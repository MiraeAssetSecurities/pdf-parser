import io
import base64
import json
import logging
import subprocess
import tempfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Optional, Union
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

import pdfplumber
from striprtf.striprtf import rtf_to_text

from .types import (
    OfficeParserConfig, OfficeParserAST, OfficeContentNode,
    OfficeMetadata, TextFormatting, OfficeAttachment
)

logger = logging.getLogger("office_parser")

_GENERIC_TITLES = {"PowerPoint Presentation", "Word Document", "Microsoft Word Document",
                   "Microsoft PowerPoint Presentation", "Presentation", "Document"}

def _clean_title(title: Optional[str]) -> Optional[str]:
    if not title or title.strip() in _GENERIC_TITLES:
        return None
    return title.strip()


def _pptx_to_slide_images(data: bytes, dpi: int = 150) -> list:
    """LibreOffice + pdf2image로 pptx를 슬라이드별 PNG 바이트로 변환.
    반환: [(slide_num, png_bytes), ...] 또는 실패 시 빈 리스트
    """
    soffice = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if not Path(soffice).exists():
        import shutil
        soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        logger.warning("⚠️ LibreOffice not found — 슬라이드 이미지 요약을 건너뜁니다")
        return []

    try:
        from pdf2image import convert_from_path
    except ImportError:
        logger.warning("⚠️ pdf2image not installed — pip install pdf2image")
        return []

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "input.pptx"
            pptx_path.write_bytes(data)
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", str(pptx_path), "--outdir", tmpdir],
                capture_output=True, timeout=120
            )
            pdf_path = Path(tmpdir) / "input.pdf"
            if not pdf_path.exists():
                logger.warning("⚠️ LibreOffice PDF 변환 실패")
                return []
            images = convert_from_path(str(pdf_path), dpi=dpi)
            result = []
            for i, img in enumerate(images, 1):
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                result.append((i, buf.getvalue()))
            return result
    except Exception as e:
        logger.warning("⚠️ 슬라이드 이미지 변환 실패: %s", e)
        return []


class OfficeParser:
    @staticmethod
    def parse_office(
        file: Union[str, bytes, Path],
        config: OfficeParserConfig = None
    ) -> OfficeParserAST:
        if config is None:
            config = OfficeParserConfig()
        
        if isinstance(file, (str, Path)):
            file_path = Path(file)
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file}")
            ext = file_path.suffix.lower().lstrip('.')
            with open(file_path, 'rb') as f:
                data = f.read()
        elif isinstance(file, bytes):
            data = file
            ext = _detect_extension(data)
        else:
            raise ValueError("Invalid file input")
        
        if ext == 'docx':
            return _parse_docx(data, config)
        elif ext == 'pptx':
            return _parse_pptx(data, config)
        elif ext == 'xlsx':
            return _parse_xlsx(data, config)
        elif ext == 'pdf':
            return _parse_pdf(data, config)
        elif ext == 'rtf':
            return _parse_rtf(data, config)
        else:
            raise ValueError(f"Unsupported file type: {ext}")


def _detect_extension(data: bytes) -> str:
    if data.startswith(b'PK'):
        return 'docx'
    elif data.startswith(b'%PDF'):
        return 'pdf'
    elif data.startswith(b'{\\rtf'):
        return 'rtf'
    raise ValueError("Cannot detect file type")


def _extract_theme_colors(wb) -> list:
    """워크북에서 테마 색상 팔레트 추출"""
    try:
        import xml.etree.ElementTree as ET
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        root = ET.fromstring(wb.loaded_theme)
        colors = []
        for scheme in root.findall('.//a:clrScheme', ns):
            for child in scheme:
                for c in child:
                    val = c.get('val', c.get('lastClr', ''))
                    colors.append(val)
        return colors
    except Exception:
        return []


def _resolve_color(color_obj, theme_colors: list) -> str:
    """openpyxl Color 객체를 #RRGGBB 문자열로 변환"""
    if not color_obj:
        return None
    try:
        if color_obj.type == 'rgb' and color_obj.rgb:
            rgb = str(color_obj.rgb)
            if rgb != '00000000':
                return '#' + rgb[2:]  # AARRGGBB -> #RRGGBB
        elif color_obj.type == 'theme' and theme_colors:
            idx = color_obj.theme
            # Excel은 테마 인덱스 0↔1, 2↔3을 교차 매핑
            if idx == 0: idx = 1
            elif idx == 1: idx = 0
            elif idx == 2: idx = 3
            elif idx == 3: idx = 2
            if 0 <= idx < len(theme_colors):
                val = theme_colors[idx]
                if val.startswith('#'):
                    return val
                elif val in ('window', 'windowText'):
                    return '#FFFFFF' if val == 'window' else '#000000'
                elif len(val) == 6:
                    return '#' + val
    except Exception:
        pass
    return None


def _extract_cell_style(cell, theme_colors: list) -> dict:
    """셀의 배경색, 글자색, 볼드 정보 추출"""
    style = {}
    try:
        if cell.fill and cell.fill.fill_type == 'solid':
            bg = _resolve_color(cell.fill.fgColor, theme_colors)
            if bg:
                style['background-color'] = bg
        if cell.font:
            fc = _resolve_color(cell.font.color, theme_colors)
            if fc:
                style['color'] = fc
            if cell.font.bold:
                style['font-weight'] = 'bold'
    except Exception:
        pass
    return style if style else None


def _extract_docx_images(element, doc, qn) -> list:
    """문단/run에서 인라인 이미지 추출. 반환: [(img_bytes, ext), ...]"""
    images = []
    for blip in element.findall(f".//{qn('a:blip')}"):
        rId = blip.get(qn("r:embed"))
        if not rId:
            continue
        try:
            rel = doc.part.rels[rId]
            img_data = rel.target_part.blob
            ct = rel.target_part.content_type or ""
            ext = ct.split("/")[-1].replace("jpeg", "jpg") if "/" in ct else "png"
            images.append((img_data, ext))
        except Exception:
            pass
    return images


def _extract_section_text(section_node) -> str:
    """section 노드에서 텍스트 추출"""
    parts = []
    if not section_node.children:
        return section_node.text or ""
    for child in section_node.children:
        if child.text:
            parts.append(child.text)
        if child.type == "table" and child.children:
            for row in child.children:
                if row.children:
                    parts.append(" | ".join(c.text or "" for c in row.children))
    return "\n".join(parts)


def _parse_docx(data: bytes, config: OfficeParserConfig) -> OfficeParserAST:
    from docx.oxml.ns import qn

    doc = Document(io.BytesIO(data))
    attachments = []
    img_counter = [0]
    # 플랫하게 모든 노드를 먼저 수집
    flat_nodes = []

    for element in doc.element.body:
        # 테이블
        if element.tag == qn("w:tbl"):
            from docx.table import Table as DocxTable
            tbl = DocxTable(element, doc)
            rows = []
            for row in tbl.rows:
                cells = [cell.text or "" for cell in row.cells]
                rows.append(OfficeContentNode(type="row", children=[
                    OfficeContentNode(type="cell", text=c) for c in cells
                ]))
            if rows:
                flat_nodes.append(OfficeContentNode(type="table", metadata={}, children=rows))
            continue

        if element.tag != qn("w:p"):
            continue

        pPr = element.find(qn("w:pPr"))

        # 이미지 추출
        if config.extract_attachments:
            for img_data, ext in _extract_docx_images(element, doc, qn):
                img_counter[0] += 1
                filename = f"image_{img_counter[0]}.{ext}"
                attachments.append(OfficeAttachment(type="image", data=img_data, filename=filename, extension=ext))
                flat_nodes.append(OfficeContentNode(
                    type="image", metadata={"filename": filename, "format": ext}
                ))

        from docx.text.paragraph import Paragraph
        para = Paragraph(element, doc)
        text = para.text.strip()
        if not text:
            continue

        style_name = para.style.name if para.style else ""

        # Heading
        if style_name.startswith("Heading"):
            try:
                level = int(style_name.split()[-1])
            except (ValueError, IndexError):
                level = 1
            flat_nodes.append(OfficeContentNode(type="heading", text=text, metadata={"level": level}))
            continue

        # List
        numPr = pPr.find(qn("w:numPr")) if pPr is not None else None
        if numPr is not None:
            ilvl_el = numPr.find(qn("w:ilvl"))
            indent_level = int(ilvl_el.get(qn("w:val"), "0")) if ilvl_el is not None else 0
            list_type = "ordered" if "Number" in style_name or "List Number" in style_name else "unordered"
            flat_nodes.append(OfficeContentNode(
                type="list", text=text, metadata={"listType": list_type, "indent_level": indent_level}
            ))
            continue

        # 일반 문단
        indent_level = 0
        if pPr is not None:
            ind = pPr.find(qn("w:ind"))
            if ind is not None:
                left = ind.get(qn("w:left"), "0")
                try:
                    indent_level = max(0, int(left) // 720)
                except ValueError:
                    pass

        flat_nodes.append(OfficeContentNode(
            type="paragraph", text=text, metadata={"indent_level": indent_level} if indent_level > 0 else None
        ))

    # ── Heading 기반 섹션 분할 ──
    # 최상위 heading(가장 작은 level)을 기준으로 섹션 분할
    heading_levels = [n.metadata["level"] for n in flat_nodes if n.type == "heading" and n.metadata]
    has_sections = len(heading_levels) >= 2  # heading이 2개 이상이어야 섹션 분할 의미 있음

    if has_sections:
        split_level = min(heading_levels)
        content = []
        current_children = []
        current_title = None
        section_num = 0

        def _flush_section():
            nonlocal current_children, current_title, section_num
            if current_children:
                section_num += 1
                meta = {"sectionNumber": section_num}
                if current_title:
                    meta["sectionTitle"] = current_title
                content.append(OfficeContentNode(
                    type="section", text="", metadata=meta, children=list(current_children)
                ))
                current_children = []
                current_title = None

        for node in flat_nodes:
            if node.type == "heading" and node.metadata.get("level") == split_level:
                _flush_section()
                current_title = node.text
                current_children.append(node)
            else:
                current_children.append(node)

        _flush_section()
    else:
        # 섹션 분할 없이 플랫하게
        content = flat_nodes

    metadata = OfficeMetadata(
        title=_clean_title(doc.core_properties.title), author=doc.core_properties.author,
        created=doc.core_properties.created, modified=doc.core_properties.modified,
    )
    ast = OfficeParserAST(type="docx", metadata=metadata, content=content,
                          attachments=attachments if attachments else None)

    # ── Bedrock 요약 ──
    if config.summarize:
        # 1) 섹션별 요약 (병렬) — 섹션이 여러 개일 때만
        sections = [n for n in content if n.type == "section"]
        if len(sections) >= 2:
            logger.info("🧠 Generating section summaries... (%d sections)", len(sections))
            with ThreadPoolExecutor() as executor:
                sec_futures = {}
                for node in sections:
                    text = _extract_section_text(node)
                    if text.strip():
                        label = node.metadata.get("sectionTitle", f"Section {node.metadata['sectionNumber']}")
                        f = executor.submit(_summarize_text, text, label, config)
                        sec_futures[f] = node
                for f in as_completed(sec_futures):
                    node = sec_futures[f]
                    try:
                        node.metadata["section_summary"] = f.result()
                        logger.info("✅ Section '%s' summary done", node.metadata.get("sectionTitle", ""))
                    except Exception as e:
                        logger.warning("⚠️ Section summary failed: %s", e)

        # 2) 이미지 요약 (병렬) — section_summary를 context로
        img_nodes_with_data = []
        def _collect_images(nodes, ctx=""):
            for node in nodes:
                if node.type == "image" and node.metadata:
                    fname = node.metadata.get("filename", "")
                    att = next((a for a in (ast.attachments or []) if a.filename == fname), None)
                    if att and _is_large_image(att.data, config.min_image_size):
                        img_nodes_with_data.append((node, att.data, att.extension, ctx))
                elif node.type == "section" and node.children:
                    sec_ctx = node.metadata.get("section_summary", "") if node.metadata else ""
                    _collect_images(node.children, sec_ctx)
                elif node.children:
                    _collect_images(node.children, ctx)
        _collect_images(content)

        if img_nodes_with_data:
            logger.info("🖼️ Generating image summaries... (%d images)", len(img_nodes_with_data))
            with ThreadPoolExecutor() as executor:
                img_futures = {}
                for child, img_data, ext, ctx in img_nodes_with_data:
                    f = executor.submit(_summarize_image, img_data, ext, config, ctx)
                    img_futures[f] = child
                for f in as_completed(img_futures):
                    child = img_futures[f]
                    try:
                        child.metadata["image_summary"] = f.result()
                        logger.info("✅ Image %s summary done", child.metadata.get("filename"))
                    except Exception as e:
                        logger.warning("⚠️ Image summary failed: %s", e)

        # 2.5) 테이블 요약 (병렬) — section_summary를 context로
        table_nodes_with_ctx = []
        def _collect_tables(nodes, ctx=""):
            for node in nodes:
                if node.type == "table" and node.children:
                    tbl_text = "\n".join(
                        " | ".join(c.text or "" for c in row.children)
                        for row in node.children if row.type == "row" and row.children
                    )
                    if tbl_text.strip():
                        table_nodes_with_ctx.append((node, tbl_text, ctx))
                elif node.type == "section" and node.children:
                    sec_ctx = node.metadata.get("section_summary", "") if node.metadata else ""
                    _collect_tables(node.children, sec_ctx)
        _collect_tables(content)

        if table_nodes_with_ctx:
            logger.info("📊 Generating table summaries... (%d tables)", len(table_nodes_with_ctx))
            with ThreadPoolExecutor() as executor:
                tbl_futures = {}
                for child, tbl_text, ctx in table_nodes_with_ctx:
                    f = executor.submit(_summarize_table, tbl_text, config, ctx)
                    tbl_futures[f] = child
                for f in as_completed(tbl_futures):
                    child = tbl_futures[f]
                    try:
                        if not child.metadata:
                            child.metadata = {}
                        child.metadata["table_summary"] = f.result()
                        logger.info("✅ Table summary done")
                    except Exception as e:
                        logger.warning("⚠️ Table summary failed: %s", e)

        # 3) 전체 문서 요약
        logger.info("📝 Generating document summary...")
        try:
            doc_text = ast.to_text()
            ast.metadata.document_summary = _summarize_document(doc_text, config, "Word document")
            logger.info("✅ Document summary done")
        except Exception as e:
            logger.warning("⚠️ Document summary failed: %s", e)

    return ast


def _pptx_shape_to_node(shape, slide_idx: int, img_counter: list, config: OfficeParserConfig, attachments: list, slide_area: int = 0):
    """단일 shape을 OfficeContentNode로 변환"""
    from pptx.shapes.group import GroupShape
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    nodes = []

    # 테이블
    if shape.has_table:
        tbl = shape.table
        rows = []
        for row in tbl.rows:
            cells = [cell.text or "" for cell in row.cells]
            rows.append(cells)
        max_cols = max((len(r) for r in rows), default=0)
        tbl_node = OfficeContentNode(type="table", metadata={"rows": len(rows), "cols": max_cols}, children=[])
        for ri, cells in enumerate(rows):
            row_node = OfficeContentNode(type="row", metadata={"row": ri}, children=[])
            for c in cells:
                row_node.children.append(OfficeContentNode(type="cell", text=c))
            tbl_node.children.append(row_node)
        nodes.append(tbl_node)

    # 이미지 — 슬라이드 면적의 30% 미만이면 스킵
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            if slide_area > 0:
                shape_area = shape.width * shape.height
                if shape_area / slide_area < 0.3:
                    return nodes  # 작은 이미지 — 추출/저장/요약 모두 스킵

            img = shape.image
            ext = img.content_type.split("/")[-1].replace("jpeg", "jpg")
            img_data = img.blob
            idx = img_counter[0]
            img_counter[0] += 1
            filename = f"slide_{slide_idx}_image_{idx}.{ext}"
            meta = {"format": ext, "slideNumber": slide_idx, "imageIndex": idx, "filename": filename}
            try:
                meta["bbox"] = {
                    "left": round(shape.left / 914400, 2),
                    "top": round(shape.top / 914400, 2),
                    "width": round(shape.width / 914400, 2),
                    "height": round(shape.height / 914400, 2),
                }
            except Exception:
                pass
            img_node = OfficeContentNode(type="image", metadata=meta)
            if config.extract_attachments and img_data:
                attachments.append(OfficeAttachment(type="image", data=img_data, filename=filename, extension=ext))
            nodes.append((img_node, img_data, ext))
        except Exception:
            pass

    # 그룹 shape — 재귀
    elif isinstance(shape, GroupShape):
        for child_shape in shape.shapes:
            nodes.extend(_pptx_shape_to_node(child_shape, slide_idx, img_counter, config, attachments, slide_area))

    # 텍스트 (text_frame 있는 shape)
    elif shape.has_text_frame:
        tf = shape.text_frame
        for para in tf.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            level = para.level or 0
            p_node = OfficeContentNode(type="paragraph", text=text, metadata={"indent_level": level})
            if para.runs:
                run = para.runs[0]
                fmt = TextFormatting(bold=run.font.bold, italic=run.font.italic)
                if fmt.bold or fmt.italic:
                    p_node.formatting = fmt
            nodes.append(p_node)

    return nodes


def _parse_pptx(data: bytes, config: OfficeParserConfig) -> OfficeParserAST:
    prs = Presentation(io.BytesIO(data))
    content = []
    attachments = []
    slide_images = {}  # slide_idx -> [(img_node, img_data, ext)]

    logger.info("📂 Parsing PowerPoint file (%d slides)", len(prs.slides))

    slide_area = prs.slide_width * prs.slide_height

    for i, slide in enumerate(prs.slides, 1):
        # 숨겨진 슬라이드 건너뛰기
        if slide._element.get('show') == '0':
            logger.info("⏭️ Slide %d (hidden) — skipped", i)
            continue
        logger.info("📄 Parsing slide %d...", i)
        slide_meta = {"slideNumber": i}

        # 슬라이드 제목 추출
        if slide.shapes.title:
            slide_meta["slideTitle"] = slide.shapes.title.text

        slide_node = OfficeContentNode(type="slide", metadata=slide_meta, children=[])
        img_counter = [0]
        slide_img_list = []

        for shape in slide.shapes:
            results = _pptx_shape_to_node(shape, i, img_counter, config, attachments, slide_area)
            for r in results:
                if isinstance(r, tuple):  # (img_node, img_data, ext)
                    img_node, img_data, ext = r
                    slide_node.children.append(img_node)
                    slide_img_list.append((img_node, img_data, ext))
                else:
                    slide_node.children.append(r)

        # 슬라이드 노트
        if not config.ignore_notes and slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_node.children.append(OfficeContentNode(type="notes", text=notes_text))

        content.append(slide_node)
        slide_images[i] = slide_img_list

    metadata = OfficeMetadata(title=_clean_title(prs.core_properties.title), author=prs.core_properties.author)

    # ── Bedrock 요약: 슬라이드 전체 이미지 기반 ──
    if config.summarize:
        # LibreOffice로 슬라이드별 이미지 생성
        logger.info("🖼️ Converting slides to images (LibreOffice)...")
        slide_pngs = _pptx_to_slide_images(data)  # [(slide_num, png_bytes), ...]
        slide_png_map = {snum: png for snum, png in slide_pngs}

        if slide_png_map:
            logger.info("🧠 Generating slide summaries from images... (%d slides, parallel)", len(slide_png_map))
            # 슬라이드 이미지 첨부파일로 저장 + 요약 병렬
            with ThreadPoolExecutor() as executor:
                futures = {}
                for node in content:
                    if node.type != "slide":
                        continue
                    snum = node.metadata["slideNumber"]
                    png_data = slide_png_map.get(snum)
                    if not png_data:
                        continue
                    # 슬라이드 전체 이미지를 첨부파일로 저장
                    filename = f"slide_{snum}.png"
                    attachments.append(OfficeAttachment(type="image", data=png_data, filename=filename, extension="png"))
                    node.metadata["slide_image"] = filename

                    slide_text = _extract_slide_text(node)
                    f = executor.submit(_summarize_slide_image, png_data, slide_text, config)
                    futures[f] = node
                for f in as_completed(futures):
                    node = futures[f]
                    try:
                        node.metadata["slide_summary"] = f.result()
                        logger.info("✅ Slide %d summary done", node.metadata["slideNumber"])
                    except Exception as e:
                        logger.warning("⚠️ Slide %d summary failed: %s", node.metadata["slideNumber"], e)
        else:
            # LibreOffice 없으면 텍스트 기반 요약으로 폴백
            logger.info("🧠 Generating slide summaries from text... (parallel)")
            with ThreadPoolExecutor() as executor:
                futures = {}
                for node in content:
                    if node.type == "slide":
                        text = _extract_slide_text(node)
                        if text.strip():
                            snum = node.metadata["slideNumber"]
                            f = executor.submit(_summarize_text, text, f"Slide {snum}", config)
                            futures[f] = node
                for f in as_completed(futures):
                    node = futures[f]
                    try:
                        node.metadata["slide_summary"] = f.result()
                        logger.info("✅ Slide %d summary done", node.metadata["slideNumber"])
                    except Exception as e:
                        logger.warning("⚠️ Slide %d summary failed: %s", node.metadata["slideNumber"], e)

    ast = OfficeParserAST(
        type="pptx", metadata=metadata, content=content,
        attachments=attachments if attachments else None
    )

    # ── 덱 전체 요약 ──
    if config.summarize:
        logger.info("📝 Generating document summary...")
        try:
            doc_text = ast.to_text()
            ast.metadata.document_summary = _summarize_document(doc_text, config, "presentation deck")
            logger.info("✅ Document summary done")
        except Exception as e:
            logger.warning("⚠️ Deck summary failed: %s", e)

    return ast


def _extract_slide_text(slide_node: OfficeContentNode) -> str:
    """슬라이드 노드에서 텍스트 추출"""
    parts = []
    if not slide_node.children:
        return ""
    for child in slide_node.children:
        if child.type == "paragraph" and child.text:
            parts.append(child.text)
        elif child.type == "table" and child.children:
            for row in child.children:
                if row.children:
                    parts.append(" | ".join(c.text or "" for c in row.children))
        elif child.type == "notes" and child.text:
            parts.append(f"[Notes] {child.text}")
    return "\n".join(parts)


def _parse_xlsx(data: bytes, config: OfficeParserConfig) -> OfficeParserAST:
    wb = load_workbook(io.BytesIO(data), data_only=True)
    content = []
    attachments = []

    logger.info("📂 Parsing Excel file (%d sheets)", len(wb.sheetnames))

    # ── 1단계: 모든 시트의 구조 파싱 (이미지 바이너리 수집) ──
    # sheet_images: { sheet_name: [(img_index_in_sheet, img_node, img_data, ext)] }
    sheet_images = {}

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        logger.info("📄 Parsing sheet '%s'...", sheet_name)

        sheet_node = OfficeContentNode(
            type="sheet",
            metadata={"sheetName": sheet_name, "maxRow": ws.max_row, "maxColumn": ws.max_column},
            children=[]
        )
        positioned = []
        sheet_img_list = []
        img_idx = 0

        # 이미지
        if hasattr(ws, '_images'):
            for img in ws._images:
                ext = getattr(img, 'format', 'png') or 'png'
                anchor_row, anchor_col = 0, 0
                try:
                    if hasattr(img.anchor, '_from'):
                        anchor_row = img.anchor._from.row
                        anchor_col = img.anchor._from.col
                except Exception:
                    pass

                img_meta = {"row": anchor_row, "col": anchor_col, "format": ext, "sheetName": sheet_name, "imageIndex": img_idx}
                img_node = OfficeContentNode(type="image", metadata=img_meta)
                positioned.append((anchor_row, img_node))

                img_data = None
                if config.extract_attachments:
                    try:
                        img_data = img._data()
                    except Exception:
                        pass

                # 시트별 일련번호로 파일명 생성
                safe_name = sheet_name.replace(' ', '_').replace('/', '_')
                filename = f"{safe_name}_image_{img_idx}.{ext}"
                if img_data and config.extract_attachments:
                    attachments.append(OfficeAttachment(
                        type="image", data=img_data, filename=filename, extension=ext
                    ))
                    img_node.metadata["filename"] = filename

                sheet_img_list.append((img_idx, img_node, img_data, ext))
                img_idx += 1

        # 차트
        if hasattr(ws, '_charts'):
            for chart in ws._charts:
                title = None
                if hasattr(chart, 'title') and chart.title:
                    try:
                        title = str(chart.title.tx.rich.p[0].r[0].t) if chart.title.tx and chart.title.tx.rich else None
                    except:
                        pass
                chart_row = 0
                try:
                    if hasattr(chart.anchor, '_from'):
                        chart_row = chart.anchor._from.row
                except Exception:
                    pass
                chart_node = OfficeContentNode(
                    type="chart", metadata={"chartType": chart.__class__.__name__, "title": title, "row": chart_row}
                )
                positioned.append((chart_row, chart_node))

        # 테마 색상 팔레트 추출
        theme_colors = _extract_theme_colors(wb)

        # 병합 셀 정보 수집: (row, col) -> colspan
        merged_spans = {}
        for mc in ws.merged_cells.ranges:
            colspan = mc.max_col - mc.min_col + 1
            if colspan > 1:
                merged_spans[(mc.min_row, mc.min_col)] = colspan
                for col in range(mc.min_col + 1, mc.max_col + 1):
                    merged_spans[(mc.min_row, col)] = 0  # 병합된 나머지 셀은 스킵

        # 셀 데이터 — 빈 셀도 위치 유지를 위해 포함하되, 뒤쪽 빈 셀은 제거
        for row_idx, row in enumerate(ws.iter_rows(), 1):
            # 마지막으로 값이 있는 셀 위치 찾기 (병합 스킵 셀 제외)
            last_val_idx = -1
            for i, cell in enumerate(row):
                span = merged_spans.get((cell.row, cell.column))
                if span == 0:
                    continue  # 병합된 나머지 셀
                if cell.value is not None:
                    last_val_idx = i
            if last_val_idx < 0:
                continue
            row_node = OfficeContentNode(type="row", metadata={"row": row_idx}, children=[])
            for cell in row[:last_val_idx + 1]:
                span = merged_spans.get((cell.row, cell.column))
                if span == 0:
                    continue  # 병합된 나머지 셀 스킵
                meta = {"row": cell.row, "col": cell.column}
                if span and span > 1:
                    meta["colspan"] = span
                # 셀 스타일 수집
                style = _extract_cell_style(cell, theme_colors)
                if style:
                    meta["style"] = style
                row_node.children.append(
                    OfficeContentNode(type="cell", text=str(cell.value) if cell.value is not None else "", metadata=meta)
                )
            positioned.append((row_idx - 1, row_node))

        positioned.sort(key=lambda x: x[0])
        sheet_node.children = [node for _, node in positioned]
        content.append(sheet_node)
        sheet_images[sheet_name] = sheet_img_list

    # ── 2단계: 시트 요약 병렬 처리 ──
    if config.summarize:
        logger.info("🧠 Generating sheet summaries... (parallel)")
        sheet_texts = {}
        for node in content:
            if node.type == "sheet":
                name = node.metadata["sheetName"]
                sheet_texts[name] = _extract_sheet_text(node)

        with ThreadPoolExecutor() as executor:
            futures = {}
            for name, text in sheet_texts.items():
                f = executor.submit(_summarize_text, text, name, config)
                futures[f] = name
            for f in as_completed(futures):
                name = futures[f]
                try:
                    summary = f.result()
                    # content에서 해당 시트 노드 찾아서 메타데이터에 요약 추가
                    for node in content:
                        if node.type == "sheet" and node.metadata["sheetName"] == name:
                            node.metadata["sheet_summary"] = summary
                            break
                    logger.info("✅ Sheet '%s' summary done", name)
                except Exception as e:
                    logger.warning("⚠️ Sheet '%s' summary failed: %s", name, e)

        # ── 3단계: 이미지 요약 병렬 처리 (시트 요약을 context로 전달) ──
        all_img_tasks = []
        for name, img_list in sheet_images.items():
            # 해당 시트의 요약 가져오기
            sheet_summary = ""
            for node in content:
                if node.type == "sheet" and node.metadata["sheetName"] == name:
                    sheet_summary = node.metadata.get("sheet_summary", "")
                    break
            for img_idx, img_node, img_data, ext in img_list:
                if img_data and _is_large_image(img_data, config.min_image_size):
                    all_img_tasks.append((name, img_node, img_data, ext, sheet_summary))
                elif img_data:
                    img_node.metadata["skipped_summary"] = "too_small"

        if all_img_tasks:
            logger.info("🖼️ Generating image summaries... (%d, parallel)", len(all_img_tasks))
            with ThreadPoolExecutor() as executor:
                futures = {}
                for name, img_node, img_data, ext, sheet_summary in all_img_tasks:
                    f = executor.submit(_summarize_image, img_data, ext, config, sheet_summary)
                    futures[f] = (name, img_node)
                for f in as_completed(futures):
                    name, img_node = futures[f]
                    try:
                        img_node.metadata["image_summary"] = f.result()
                        logger.info("✅ Image summary done (sheet: '%s', row: %s)", name, img_node.metadata.get("row"))
                    except Exception as e:
                        logger.warning("⚠️ Image summary failed (sheet: '%s'): %s", name, e)

    return OfficeParserAST(
        type="xlsx", metadata=OfficeMetadata(), content=content,
        attachments=attachments if attachments else None
    )


def _extract_sheet_text(sheet_node: OfficeContentNode) -> str:
    texts = []
    for child in sheet_node.children:
        if child.type == "row" and child.children:
            texts.append(" | ".join(c.text for c in child.children if c.text))
        elif child.type == "chart" and child.metadata:
            texts.append(f"[차트: {child.metadata.get('chartType', '')}] {child.metadata.get('title', '')}")
    return "\n".join(texts)


_bedrock_client_cache = {}

def _get_bedrock_client(config: OfficeParserConfig):
    import boto3
    region = config.bedrock_region
    if region not in _bedrock_client_cache:
        _bedrock_client_cache[region] = boto3.client("bedrock-runtime", region_name=region)
    return _bedrock_client_cache[region]


def _is_large_image(img_data: bytes, min_size: int) -> bool:
    """이미지가 min_size x min_size 이상인지 확인"""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(img_data))
        return img.width >= min_size and img.height >= min_size
    except Exception:
        return False


def _summarize_image(img_data: bytes, fmt: str, config: OfficeParserConfig, sheet_summary: str = "") -> str:
    client = _get_bedrock_client(config)
    media_type = f"image/{fmt}" if fmt != "jpg" else "image/jpeg"
    b64 = base64.b64encode(img_data).decode("utf-8")

    prompt = "Describe this image in detail in 3-5 sentences. Write in Korean."
    if sheet_summary:
        prompt = f"Context of the sheet containing this image: {sheet_summary}\n\nUsing the above context, describe this image in detail in 3-5 sentences. Write in Korean."

    body = json.dumps({
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 500,
        "messages": [{"role": "user", "content": [
            {"type": "image", "source": {"type": "base64", "media_type": media_type, "data": b64}},
            {"type": "text", "text": prompt}
        ]}]
    })
    resp = client.invoke_model(modelId=config.bedrock_model_id, body=body)
    return json.loads(resp["body"].read())["content"][0]["text"]


def _summarize_slide_image(img_data: bytes, slide_text: str, config: OfficeParserConfig) -> str:
    """슬라이드 전체 이미지 + 텍스트 컨텍스트로 요약"""
    client = _get_bedrock_client(config)
    b64 = base64.b64encode(img_data).decode("utf-8")
    prompt = "Summarize this presentation slide in 5-8 sentences. Cover the key points, and if there are diagrams, charts, or architecture figures, explain their meaning and relationships. Write in Korean."
    if slide_text:
        prompt += f"\n\nSlide text:\n{slide_text[:3000]}"
    body = json.dumps({
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 800,
        "messages": [{"role": "user", "content": [
            {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": b64}},
            {"type": "text", "text": prompt}
        ]}]
    })
    resp = client.invoke_model(modelId=config.bedrock_model_id, body=body)
    return json.loads(resp["body"].read())["content"][0]["text"]



_DECK_SUMMARY_CHUNK_SIZE = 30000  # Bedrock 입력 제한 고려

def _summarize_document(full_text: str, config: OfficeParserConfig, doc_type: str = "document") -> str:
    """전체 문서 텍스트를 청킹하여 요약 후 최종 요약 생성"""
    client = _get_bedrock_client(config)

    if len(full_text) <= _DECK_SUMMARY_CHUNK_SIZE:
        chunks = [full_text]
    else:
        chunks = [full_text[i:i + _DECK_SUMMARY_CHUNK_SIZE]
                  for i in range(0, len(full_text), _DECK_SUMMARY_CHUNK_SIZE)]

    summaries = []
    for i, chunk in enumerate(chunks):
        prompt = (f"Below is part {i+1}/{len(chunks)} of a {doc_type}.\n"
                  f"Summarize the key points in 5-8 sentences. Write in Korean.\n\n{chunk}")
        body = json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 1000,
            "messages": [{"role": "user", "content": prompt}]
        })
        resp = client.invoke_model(modelId=config.bedrock_model_id, body=body)
        summaries.append(json.loads(resp["body"].read())["content"][0]["text"])

    if len(summaries) == 1:
        return summaries[0]

    combined = "\n\n".join(summaries)
    prompt = (f"Below are partial summaries of a {doc_type}.\n"
              "Create a final comprehensive summary in 5-10 sentences. Write in Korean.\n\n" + combined)
    body = json.dumps({
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 1000,
        "messages": [{"role": "user", "content": prompt}]
    })
    resp = client.invoke_model(modelId=config.bedrock_model_id, body=body)
    return json.loads(resp["body"].read())["content"][0]["text"]

def _summarize_text(text: str, sheet_name: str, config: OfficeParserConfig) -> str:
    client = _get_bedrock_client(config)
    truncated = text[:4000]
    body = json.dumps({
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 500,
        "messages": [{"role": "user",
                      "content": f"Below is the data from '{sheet_name}'. Summarize the content in 3-5 sentences. Write in Korean.\n\n{truncated}"}]
    })
    resp = client.invoke_model(modelId=config.bedrock_model_id, body=body)
    return json.loads(resp["body"].read())["content"][0]["text"]


def _summarize_table(table_text: str, config: OfficeParserConfig, context: str = "") -> str:
    client = _get_bedrock_client(config)
    prompt = "Summarize this table in 2-4 sentences. Describe what data it contains and key insights. Write in Korean."
    if context:
        prompt = f"Context of the section containing this table:\n{context[:2000]}\n\n{prompt}"
    prompt += f"\n\nTable:\n{table_text[:4000]}"
    body = json.dumps({
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 500,
        "messages": [{"role": "user", "content": prompt}]
    })
    resp = client.invoke_model(modelId=config.bedrock_model_id, body=body)
    return json.loads(resp["body"].read())["content"][0]["text"]


def _parse_pdf(data: bytes, config: OfficeParserConfig) -> OfficeParserAST:
    content = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text_content = page.extract_text() or ""
            if text_content.strip():
                content.append(OfficeContentNode(type="page", text=text_content, metadata={"pageNumber": i}))
        metadata = OfficeMetadata(pages=len(pdf.pages))
    return OfficeParserAST(type="pdf", metadata=metadata, content=content)


def _parse_rtf(data: bytes, config: OfficeParserConfig) -> OfficeParserAST:
    text_content = rtf_to_text(data.decode('utf-8', errors='ignore'))
    content = [OfficeContentNode(type="paragraph", text=text_content)]
    return OfficeParserAST(type="rtf", metadata=OfficeMetadata(), content=content)
